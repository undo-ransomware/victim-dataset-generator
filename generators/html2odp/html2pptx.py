import requests
from lxml import etree
import re
from collections import Counter
import nltk
from nltk.corpus import stopwords
from mwclient import Site
import os
from slugify import slugify
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image
from math import sqrt
from resizeimage import resizeimage

stop = set(stopwords.words('english'))
nopunct = re.compile('\w')
site = Site('en.wikipedia.org')

def download_nltk():
	try:
		stopwords.words("english")
	except LookupError:
		import nltk
		nltk.download('stopwords')
		nltk.download('punkt')

def tokenize(text):
	return [word for word in nltk.word_tokenize(text.lower()) if word not in stop and nopunct.match(word)]

def extract_keywords(text):
	return dict(Counter(tokenize(text)).most_common(20))

def score(sent, keywords):
	words = tokenize(sent)
	if len(words) == 0:
		return -1
	return sum(keywords[word] for word in words if word in keywords) / (len(words) + 3.0)

def summarize(text, keywords, limit):
	scores = dict()
	sentences = nltk.sent_tokenize(text)
	if len(sentences) == 0:
		return []
	for sent in sentences:
		scores[sent] = score(sent, keywords)
	max_score = max(scores.values())

	important = set()
	for sent in sorted(sentences, key=lambda s: scores[s], reverse=True):
		if scores[sent] < max_score / 2.5:
			break;
		important.add(sent)
		if len(important) == 7:
			break;
		if len(" ".join(important)) > limit:
			break;
	return [sent for sent in sentences if sent in important]

def make_title(prs, title, subtitle):
	slide = prs.slides.add_slide(prs.slide_layouts[0])
	slide.shapes.title.text = title
	slide.placeholders[1].text = subtitle

def make_bulleted(prs, title, bullets):
	slide = prs.slides.add_slide(prs.slide_layouts[1])
	slide.shapes.title.text = title
	textbox = slide.shapes.placeholders[1].text_frame
	textbox.text = "\n".join(bullets)
	# in classic bad presentation style, reduce font size until the text fits.
	# 240 is a rough estimate for how much text fits on one page at size 36.
	size = max(min(sqrt(240) / sqrt(len(textbox.text)) * 36, 36), 1)
	for p in textbox.paragraphs:
		p.font.size = Pt(size)

def make_image(prs, filename, caption):
	with Image.open(filename) as img:
		width, height = img.size

	# usable size is 5.5" x 8.0" in the idiotic-inch based default template.
	# in classic bad presentation style, we use most of that for the caption.
	WIDTH, HEIGHT, MARGIN = 8.0, 5.0, 1.0
	w = h = None
	x = y = Inches(MARGIN)
	if height / HEIGHT > width / WIDTH:
		h = Inches(HEIGHT)
		x = Inches((WIDTH - HEIGHT / height * width) / 2 + MARGIN)
	else:
		w = Inches(WIDTH)
		y = Inches((HEIGHT - WIDTH / width * height) / 2 + MARGIN)
	slide = prs.slides.add_slide(prs.slide_layouts[6])
	slide.shapes.add_picture(filename, x, y, width=w, height=h)

	x, y = Inches(MARGIN), Inches(HEIGHT + 1.2 * MARGIN)
	w, h = Inches(WIDTH), Inches(MARGIN)
	textbox = slide.shapes.add_textbox(x, y, w, h).text_frame
	textbox.text = caption
	# make the text fit. single-line, so roughly linear
	size = max(min(40.0 / len(textbox.text) * 36, 36), 1)
	for p in textbox.paragraphs:
		p.font.size = Pt(size)

def sectiontext(elem):
	return "".join(elem.xpath("*[not(@class) or @class!='mw-ref']//text()|text()"))

def convert(html, prs):
	xml = etree.fromstring(html)
	title = " ".join(xml.xpath("//title//text()"))
	make_title(prs, title, "From Wikipedia, the free encyclopedia");
	keywords = extract_keywords(" ".join(xml.xpath("//section/p//text()")))
	
	for section in xml.xpath("//section"):
		headers = section.xpath("*[starts-with(name(),'h')]//text()")
		if len(headers) > 0:
			header = " ".join(headers)
		else:
			header = title
		
		for figure in section.xpath('figure'):
			caption = sectiontext(figure.xpath(".//figcaption")[0])
			src = re.sub('^.*?File:', '', figure.xpath(".//a//img/@resource")[0])

			basename, ext = os.path.splitext(src)
			localname = slugify(basename, max_length=150) + ext
			if not os.path.isfile(localname):
				file = site.images[src]
				with open(localname, 'wb') as fd:
					file.download(fd)
			with Image.open(localname) as img:
				resizeimage.resize_thumbnail(img, [1200, 750]).save("temp" + ext)
			summary = "".join(summarize(caption, keywords, 60))
			make_image(prs, "temp" + ext, summary)
			os.remove("temp" + ext)
		
		text = " ".join(sectiontext(p) for p in section.xpath("p"))
		summary = summarize(text, keywords, 300)
		if len(summary) > 0:
			make_bulleted(prs, header, summary)

def download(page_title):
	prs = Presentation()
	m = requests.get("https://en.wikipedia.org/api/rest_v1/page/title/" + page_title)
	print m.text
	r = requests.get("https://en.wikipedia.org/api/rest_v1/page/html/" + page_title)
	convert(r.text, prs)
	prs.save(page_title + ".pptx")

def main(args):
	download_nltk()
	download(args[0])
	return 0
 
if __name__ == "__main__":
    import sys 
    status = main(sys.argv[1:])
    sys.exit(status)
