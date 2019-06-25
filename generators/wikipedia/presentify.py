import os
import pptx
from pptx.util import Inches, Pt
from PIL import Image
from math import sqrt
from summarizer import extract_keywords, summarize

class Presentation:
	def __init__(self):
		self.pptx = pptx.Presentation()
	
	def init_full_text(self, text):
		self.keywords = extract_keywords(text)
	
	def make_title(self, title, subtitle, *subsubtitles):
		slide = self.pptx.slides.add_slide(self.pptx.slide_layouts[0])
		slide.shapes.title.text = title
		textbox = slide.placeholders[1].text_frame
		textbox.text = subtitle
		textbox.paragraphs[0].font.size = Pt(28)
		for sst in subsubtitles:
			p = textbox.add_paragraph()
			p.text = sst
			p.font.size = Pt(16)

	def make_section(self, title, texts):
		# in classic bad presentation style, aim for somewhat more content per
		# slide than what actually fits
		summary = summarize(texts, self.keywords, 300)
		if len(summary) == 0:
			return

		slide = self.pptx.slides.add_slide(self.pptx.slide_layouts[1])
		slide.shapes.title.text = title
		textbox = slide.shapes.placeholders[1].text_frame
		textbox.text = "\n".join(summary)
		# keeping with bad presentation style, just reduce the font size until
		# all the text fits. 200 is a rough estimate for how much text fits on
		# one page at size 36.
		size = max(min(sqrt(200) / sqrt(len(textbox.text)) * 36, 36), 1)
		for p in textbox.paragraphs:
			p.font.size = Pt(size)

	def make_image(self, filename, captions):
		with Image.open(filename) as img:
			width, height = img.size

		# usable size is 8.0" x 5.5" in the idiotic-inch based default template.
		# in classic bad presentation style, we use most of that for the caption.
		WIDTH, HEIGHT, MARGIN = 8.0, 5.0, 1.0
		w = h = None
		x = y = Inches(MARGIN)
		if height / HEIGHT > width / WIDTH:
			h = Inches(HEIGHT)
			x = Inches((WIDTH - HEIGHT / height * width) / 2 + MARGIN)
		else:
			w = Inches(WIDTH)
			y = Inches((HEIGHT - WIDTH / width * height) / 2 + MARGIN)
		slide = self.pptx.slides.add_slide(self.pptx.slide_layouts[6])
		slide.shapes.add_picture(filename, x, y, width=w, height=h)
		
		summary = summarize(captions, self.keywords, 60)
		if len(summary) == 0:
			return
		x, y = Inches(MARGIN), Inches(HEIGHT + 1.2 * MARGIN)
		w, h = Inches(WIDTH), Inches(MARGIN)
		textbox = slide.shapes.add_textbox(x, y, w, h).text_frame
		textbox.text = " ".join(summary)
		# make the text fit. single-line, so roughly linear
		size = max(min(40.0 / len(textbox.text) * 36, 36), 1)
		for p in textbox.paragraphs:
			p.font.size = Pt(size)

	def save(self, filename):
		self.pptx.save('media/' + filename + '.pptx')
